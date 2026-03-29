import pdfplumber
from pptx import Presentation
import os
import sys

def extract_pdf_text(pdf_path):
    """Extract text from PDF file"""
    text = ""
    try:
        with pdfplumber.open(pdf_path) as pdf:
            for page_num, page in enumerate(pdf.pages, 1):
                page_text = page.extract_text()
                if page_text:
                    text += f"\n--- Page {page_num} ---\n{page_text}\n"
    except Exception as e:
        text = f"Error extracting PDF: {str(e)}"
    return text

def extract_pptx_text(pptx_path):
    """Extract text from PowerPoint file"""
    text = ""
    try:
        prs = Presentation(pptx_path)
        for slide_num, slide in enumerate(prs.slides, 1):
            text += f"\n--- Slide {slide_num} ---\n"
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    if shape.text.strip():
                        text += shape.text + "\n"
    except Exception as e:
        text = f"Error extracting PPTX: {str(e)}"
    return text

def main():
    notes_dir = r"C:\Users\sethc\OneDrive\Seth\Source\Money\Inventions\Programming\HTML, CSS, JavaScript\G2GPT\notes"
    output_dir = r"C:\Users\sethc\OneDrive\Seth\Source\Money\Inventions\Programming\HTML, CSS, JavaScript\G2GPT\notes"
    
    # Get all PDF and PPTX files
    files = os.listdir(notes_dir)
    
    for filename in sorted(files):
        file_path = os.path.join(notes_dir, filename)
        
        if filename.endswith('.pdf'):
            print(f"Extracting {filename}...")
            text = extract_pdf_text(file_path)
            output_file = os.path.join(output_dir, filename.replace('.pdf', '_extracted.txt'))
            with open(output_file, 'w', encoding='utf-8') as f:
                f.write(f"=== {filename} ===\n{text}")
            print(f"  -> Saved to {filename.replace('.pdf', '_extracted.txt')}")
            
        elif filename.endswith('.pptx'):
            print(f"Extracting {filename}...")
            text = extract_pptx_text(file_path)
            output_file = os.path.join(output_dir, filename.replace('.pptx', '_extracted.txt'))
            with open(output_file, 'w', encoding='utf-8') as f:
                f.write(f"=== {filename} ===\n{text}")
            print(f"  -> Saved to {filename.replace('.pptx', '_extracted.txt')}")
    
    print("\nExtraction complete!")

if __name__ == "__main__":
    main()
